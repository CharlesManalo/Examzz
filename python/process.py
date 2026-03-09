from http.server import BaseHTTPRequestHandler
import json
import io
import re
import random
import traceback
import tempfile
import os
import sys

# Add the api/src directory to Python path to import our utility
sys.path.append(os.path.join(os.path.dirname(__file__), '..', 'api', 'src'))

try:
    from utils.extract_text import extract_text_from_file, validate_file_before_extraction
    HAS_NEW_EXTRACTOR = True
except ImportError as e:
    print(f"Warning: Could not import new extractor: {e}")
    HAS_NEW_EXTRACTOR = False

# ── Text extraction (using new utility) ─────────────────────────────────────

def extract_text(filename: str, data: bytes) -> str:
    """
    Extract text from uploaded file using the new utility.
    Falls back to old method if new utility unavailable.
    """
    if HAS_NEW_EXTRACTOR:
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as temp_file:
                temp_file.write(data)
                temp_file_path = temp_file.name
            
            try:
                # Validate file
                validate_file_before_extraction(temp_file_path)
                
                # Extract text using new utility
                text = extract_text_from_file(temp_file_path)
                return text
                
            finally:
                # Clean up temporary file
                os.unlink(temp_file_path)
                
        except Exception as e:
            print(f"New extractor failed: {e}")
            # Fall back to old method
            return extract_text_fallback(filename, data)
    else:
        # Use fallback method
        return extract_text_fallback(filename, data)

def extract_text_fallback(filename: str, data: bytes) -> str:
    """
    Fallback extraction method (original implementation).
    """
    ext = filename.rsplit(".", 1)[-1].lower()
    
    # Import libraries as needed
    try:
        if ext == "pdf":
            import fitz
            doc = fitz.open(stream=data, filetype="pdf")
            pages = []
            for page in doc:
                pages.append(page.get_text())
            return "\n".join(pages)
            
        elif ext in ("docx", "doc"):
            from docx import Document
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            
        elif ext in ("pptx", "ppt"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            return "\n".join(lines)
            
        elif ext in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True)
            lines = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_text = " ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        lines.append(row_text.strip())
            return "\n".join(lines)
            
        elif ext == "txt":
            return data.decode("utf-8", errors="ignore")
            
        else:
            raise ValueError(f"Unsupported file type: {ext}")
            
    except ImportError as e:
        raise RuntimeError(f"Missing library for {ext}: {e}")

# ── Content analysis (unchanged) ─────────────────────────────────────────────

def extract_sentences(text: str) -> list[str]:
    sentences = []
    for line in text.split("\n"):
        line = line.strip()
        if not line or len(line) < 20:
            continue
        if len(line) <= 500:
            sentences.append(line)
        else:
            # Split long lines on sentence boundaries
            start = 0
            for i, ch in enumerate(line):
                if ch in ".!?" and i - start > 20:
                    chunk = line[start:i + 1].strip()
                    if 20 < len(chunk) <= 500:
                        sentences.append(chunk)
                    start = i + 1
            remainder = line[start:].strip()
            if len(remainder) > 20:
                sentences.append(remainder[:500])
    return list(dict.fromkeys(sentences))  # deduplicate, preserve order


def extract_keywords(text: str, headings: list[str]) -> list[str]:
    keywords = set()

    # Words from headings
    for h in headings:
        for w in h.split():
            if len(w) > 3:
                keywords.add(w.lower())

    # Frequent capitalised words (proper nouns / terms)
    word_freq: dict[str, int] = {}
    for w in re.findall(r'\b[A-Z][a-z]{3,}\b', text):
        word_freq[w.lower()] = word_freq.get(w.lower(), 0) + 1

    for word, count in word_freq.items():
        if count > 1:
            keywords.add(word)

    return list(keywords)[:50]


def extract_headings(text: str) -> list[str]:
    headings = []
    for line in text.split("\n"):
        t = line.strip()
        if not t or len(t) < 5 or len(t) > 100:
            continue
        if t == t.upper() or (t[0] == t[0].upper() and not t[-1] in ".!?"):
            headings.append(t)
    return list(dict.fromkeys(headings))[:20]


# ── Quiz generation (unchanged) ───────────────────────────────────────────────

def generate_mcq_questions(sentences: list[str], keywords: list[str], count: int = 5) -> list[dict]:
    questions = []
    
    for sentence in random.sample(sentences, min(count, len(sentences))):
        # Find a keyword in the sentence
        sentence_words = sentence.lower().split()
        available_keywords = [kw for kw in keywords if kw.lower() in sentence_words]
        
        if not available_keywords:
            continue
            
        keyword = random.choice(available_keywords)
        
        # Create question
        question_text = sentence.replace(keyword, "________")
        
        # Generate options
        other_keywords = [kw for kw in keywords if kw != keyword]
        distractors = random.sample(other_keywords, min(3, len(other_keywords)))
        
        if len(distractors) < 3:
            continue
            
        options = [keyword] + distractors
        random.shuffle(options)
        
        correct_index = options.index(keyword)
        
        questions.append({
            "question": question_text,
            "options": options,
            "correct_answer": correct_index,
            "type": "multiple_choice"
        })
    
    return questions


def generate_fill_blank_questions(sentences: list[str], keywords: list[str], count: int = 3) -> list[dict]:
    questions = []
    
    for sentence in random.sample(sentences, min(count, len(sentences))):
        sentence_words = sentence.lower().split()
        available_keywords = [kw for kw in keywords if kw.lower() in sentence_words and len(kw) > 4]
        
        if not available_keywords:
            continue
            
        keyword = random.choice(available_keywords)
        
        # Create question
        question_text = f"Fill in the blank: {sentence.replace(keyword, '________')}"
        
        questions.append({
            "question": question_text,
            "options": [keyword],
            "correct_answer": 0,
            "type": "fill_blank"
        })
    
    return questions


# ── HTTP Handler ─────────────────────────────────────────────────────────────

class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        content_length = int(self.headers['Content-Length'])
        post_data = self.rfile.read(content_length)
        
        try:
            # Parse multipart form data
            boundary = self.headers['Content-Type'].split('boundary=')[1]
            parts = post_data.split(f'--{boundary}'.encode())
            
            file_data = None
            filename = None
            question_count = 10
            
            for part in parts:
                if b'Content-Disposition: form-data' in part and b'filename=' in part:
                    headers_end = part.find(b'\r\n\r\n')
                    if headers_end != -1:
                        file_data = part[headers_end + 4:].rstrip(b'\r\n')
                        
                        # Extract filename
                        filename_match = re.search(rb'filename="([^"]*)"', part)
                        if filename_match:
                            filename = filename_match.group(1).decode('utf-8')
                
                elif b'name="question_count"' in part:
                    headers_end = part.find(b'\r\n\r\n')
                    if headers_end != -1:
                        question_count = int(part[headers_end + 4:].strip())
            
            if not file_data or not filename:
                raise ValueError("No file uploaded")
            
            # Extract text
            text = extract_text(filename, file_data)
            
            # Analyze content
            sentences = extract_sentences(text)
            keywords = extract_keywords(text, [])
            headings = extract_headings(text)
            
            # Generate questions
            mcq_count = min(int(question_count * 0.6), len(sentences))
            fill_count = question_count - mcq_count
            
            mcq_questions = generate_mcq_questions(sentences, keywords, mcq_count)
            fill_questions = generate_fill_blank_questions(sentences, keywords, fill_count)
            
            all_questions = mcq_questions + fill_questions
            random.shuffle(all_questions)
            
            response_data = {
                "success": True,
                "extracted_text": text[:1000] + "..." if len(text) > 1000 else text,
                "questions": all_questions[:question_count],
                "stats": {
                    "sentences_found": len(sentences),
                    "keywords_found": len(keywords),
                    "headings_found": len(headings),
                    "total_questions": len(all_questions)
                }
            }
            
        except Exception as e:
            response_data = {
                "success": False,
                "error": str(e),
                "traceback": traceback.format_exc()
            }
        
        # Send response
        self.send_response(200)
        self.send_header('Content-type', 'application/json')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        
        response = json.dumps(response_data, indent=2)
        self.wfile.write(response.encode())
