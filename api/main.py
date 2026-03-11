import os
import sys
import json
import logging
import time
import re
import hmac
import hashlib
import urllib.request
from http.server import BaseHTTPRequestHandler
from urllib.parse import urlparse, parse_qs

sys.path.insert(0, os.path.dirname(__file__))

from dotenv import load_dotenv
load_dotenv()

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ── Providers (tried in order) ────────────────────────────────────────────────
PROVIDERS = [
    {"name": "Gemini", "api_key_env": "GEMINI_API_KEY", "type": "gemini", "model": "gemini-2.0-flash"},
    {"name": "Gemini-1.5", "api_key_env": "GEMINI_API_KEY", "type": "gemini", "model": "gemini-1.5-flash"},
    {"name": "Groq", "api_key_env": "GROQ_API_KEY", "type": "openai_compat", "base_url": "https://api.groq.com/openai/v1", "model": "llama-3.3-70b-versatile"},
    {"name": "OpenRouter", "api_key_env": "OPENROUTER_API_KEY", "type": "openai_compat", "base_url": "https://openrouter.ai/api/v1", "model": "mistralai/mistral-7b-instruct:free"},
    {"name": "Hugging Face", "api_key_env": "HF_TOKEN", "type": "huggingface", "model": "Qwen/Qwen2.5-7B-Instruct"},
]


def build_prompt(extracted_text: str, question_count: int, difficulty: str) -> str:
    difficulty_prompts = {
        "easy": "simple, straightforward questions that test basic understanding",
        "medium": "moderately challenging questions that require some thinking",
        "hard": "complex questions that require deep understanding and analysis",
    }
    return f"""
Act as a professional educator. Based on the following lesson text, generate a {question_count}-question multiple choice quiz.

Requirements:
- Create {difficulty_prompts.get(difficulty, "moderately challenging")} questions
- Each question must have exactly 4 options (A, B, C, D)
- Only one option should be correct
- ONLY ask questions about LESSON CONTENT (concepts, definitions, facts, processes)
- Do NOT ask questions about the module structure, icons, how to use the module, or learning objectives
- Questions should test understanding of the actual subject matter taught
- Avoid true/false questions
- Make questions clear and unambiguous

Return the output strictly as a JSON array with this exact structure:
[
    {{
        "question": "The question text here",
        "options": ["Option A", "Option B", "Option C", "Option D"],
        "answer_index": 0,
        "explanation": "Brief explanation of why this answer is correct"
    }}
]

Do not include any text outside the JSON array. The response must be valid JSON.

Lesson text:
{extracted_text}
"""


def call_gemini(provider: dict, prompt: str) -> str:
    from google import genai
    client = genai.Client(api_key=provider["api_key"])
    response = client.models.generate_content(model=provider["model"], contents=prompt)
    return response.text


def call_huggingface(provider: dict, prompt: str) -> str:
    from huggingface_hub import InferenceClient
    client = InferenceClient(api_key=provider["api_key"])
    response = client.chat.completions.create(
        model=provider["model"],
        messages=[
            {"role": "system", "content": "You are a professional educator. Always respond with valid JSON only, no markdown."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=2048,
        stream=False
    )
    return response.choices[0].message.content


def call_openai_compat(provider: dict, prompt: str) -> str:
    payload = json.dumps({
        "model": provider["model"],
        "messages": [
            {"role": "system", "content": "You are a professional educator. Always respond with valid JSON only, no markdown."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.7,
        "max_tokens": 2048,
    }).encode()
    req = urllib.request.Request(
        f"{provider['base_url']}/chat/completions",
        data=payload,
        headers={"Content-Type": "application/json", "Authorization": f"Bearer {provider['api_key']}"},
        method="POST"
    )
    with urllib.request.urlopen(req, timeout=55) as resp:
        data = json.loads(resp.read().decode())
        return data["choices"][0]["message"]["content"]


def extract_retry_delay(error_str: str) -> int:
    match = re.search(r'retry in (\d+)', error_str, re.IGNORECASE)
    return int(match.group(1)) if match else 0


def extract_text_from_file(content: bytes, filename: str) -> str:
    fname = filename.lower()
    if fname.endswith(".pdf"):
        try:
            import pypdf, io
            reader = pypdf.PdfReader(io.BytesIO(content))
            return "".join(page.extract_text() or "" for page in reader.pages)
        except Exception as e:
            logger.warning(f"PDF extraction failed: {e}")
            return content.decode("utf-8", errors="ignore")
    elif fname.endswith(".docx"):
        try:
            import docx, io
            doc = docx.Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            logger.warning(f"DOCX extraction failed: {e}")
            return content.decode("utf-8", errors="ignore")
    elif fname.endswith(".pptx"):
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.warning(f"PPTX extraction failed: {e}")
            return content.decode("utf-8", errors="ignore")
    else:
        return content.decode("utf-8", errors="ignore")


def generate_with_fallback(prompt: str) -> tuple:
    last_error = None
    for provider in PROVIDERS:
        api_key = os.getenv(provider["api_key_env"])
        if not api_key:
            continue
        for attempt in range(2):
            try:
                p = {**provider, "api_key": api_key}
                if provider["type"] == "gemini":
                    text = call_gemini(p, prompt)
                elif provider["type"] == "huggingface":
                    text = call_huggingface(p, prompt)
                else:
                    text = call_openai_compat(p, prompt)
                logger.info(f"Success with provider: {provider['name']}")
                return text, provider["name"]
            except Exception as e:
                err_str = str(e)
                is_rate_limit = "429" in err_str
                retry_delay = extract_retry_delay(err_str)
                if is_rate_limit and attempt == 0:
                    if retry_delay > 10:
                        last_error = e
                        break
                    else:
                        time.sleep(max(retry_delay, 3))
                else:
                    last_error = e
                    break
    raise ValueError(f"All providers exhausted. Last error: {last_error}")


def generate_quiz_logic(content: bytes, filename: str, question_count: int = 10, difficulty: str = "medium"):
    extracted_text = extract_text_from_file(content, filename)
    if not extracted_text.strip():
        raise ValueError("No text extracted from file.")

    skip_phrases = ["what i need to know", "introductory message", "for the facilitator",
                    "for the learner", "how to use this module", "notes to the teacher"]
    lines = extracted_text.split("\n")
    filtered_lines, skip = [], False
    for line in lines:
        lower = line.lower().strip()
        if any(phrase in lower for phrase in skip_phrases):
            skip = True
        if "what is it" in lower or "lesson 1" in lower or "lesson 2" in lower:
            skip = False
        if not skip:
            filtered_lines.append(line)

    cleaned_text = "\n".join(filtered_lines).strip()
    final_text = cleaned_text if len(cleaned_text) > 500 else extracted_text
    if len(final_text) > 50000:
        final_text = final_text[:50000] + "..."

    prompt = build_prompt(final_text, question_count, difficulty)
    response_text, provider_used = generate_with_fallback(prompt)
    if not response_text:
        raise ValueError("Empty response from AI provider.")

    clean = response_text.strip()
    if clean.startswith("```"):
        clean = clean.split("\n", 1)[-1]
        clean = clean.rsplit("```", 1)[0].strip()

    quiz_data = json.loads(clean)
    if not isinstance(quiz_data, list):
        raise ValueError("Response is not a JSON array")

    for i, question in enumerate(quiz_data):
        for key in ["question", "options", "answer_index", "explanation"]:
            if key not in question:
                raise ValueError(f"Question {i} missing key: {key}")
        if not isinstance(question["options"], list) or len(question["options"]) != 4:
            raise ValueError(f"Question {i} must have exactly 4 options")
        if not isinstance(question["answer_index"], int) or not 0 <= question["answer_index"] <= 3:
            raise ValueError(f"Question {i} has invalid answer_index")

    return {
        "success": True,
        "quiz": quiz_data,
        "metadata": {
            "file_name": filename,
            "file_size": len(content),
            "extracted_chars": len(extracted_text),
            "question_count": len(quiz_data),
            "difficulty": difficulty,
            "provider_used": provider_used,
        }
    }


# ── PayMongo helpers ──────────────────────────────────────────────────────────

PAYMONGO_SECRET_KEY = os.getenv("PAYMONGO_SECRET_KEY", "")
PAYMONGO_WEBHOOK_SECRET = os.getenv("PAYMONGO_WEBHOOK_SECRET", "")
SUPABASE_URL = os.getenv("SUPABASE_URL") or os.getenv("VITE_SUPABASE_URL", "")
SUPABASE_SERVICE_ROLE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY", "")

SUPPORTER_AMOUNT = 10000  # ₱100.00 in centavos


def paymongo_auth_header() -> str:
    import base64
    return f"Basic {base64.b64encode(f'{PAYMONGO_SECRET_KEY}:'.encode()).decode()}"


def supabase_request(method: str, path: str, body: dict = None):
    url = f"{SUPABASE_URL}/rest/v1{path}"
    data = json.dumps(body).encode() if body else None
    req = urllib.request.Request(
        url, data=data,
        headers={
            "apikey": SUPABASE_SERVICE_ROLE_KEY,
            "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
            "Content-Type": "application/json",
            "Prefer": "return=representation",
        },
        method=method
    )
    try:
        with urllib.request.urlopen(req, timeout=10) as resp:
            return json.loads(resp.read().decode())
    except urllib.error.HTTPError as e:
        logger.error(f"Supabase error {e.code}: {e.read().decode()}")
        raise


def create_paymongo_link(user_id: str, user_email: str) -> dict:
    payload = json.dumps({
        "data": {
            "attributes": {
                "amount": SUPPORTER_AMOUNT,
                "currency": "PHP",
                "description": "EXAMZZ Supporter — Lifetime Premium Access",
                "remarks": f"user_id:{user_id}",
                "payment_method_types": ["gcash", "paymaya", "qrph", "grab_pay"],
                "metadata": {
                    "user_id": user_id,
                    "user_email": user_email,
                    "plan": "supporter",
                },
            }
        }
    }).encode()

    req = urllib.request.Request(
        "https://api.paymongo.com/v1/links",
        data=payload,
        headers={"Authorization": paymongo_auth_header(), "Content-Type": "application/json"},
        method="POST"
    )
    with urllib.request.urlopen(req, timeout=15) as resp:
        data = json.loads(resp.read().decode())
        link = data["data"]
        return {
            "link_id": link["id"],
            "checkout_url": link["attributes"]["checkout_url"],
            "reference_number": link["attributes"]["reference_number"],
        }


def verify_webhook_signature(body: bytes, sig_header: str) -> bool:
    if not PAYMONGO_WEBHOOK_SECRET:
        return True
    try:
        expected = hmac.new(PAYMONGO_WEBHOOK_SECRET.encode(), body, hashlib.sha256).hexdigest()
        return hmac.compare_digest(expected, sig_header)
    except Exception:
        return False


def mark_user_premium(user_id: str, payment_id: str, link_id: str):
    supabase_request("PATCH", f"/users?id=eq.{user_id}", {
        "is_premium": True,
        "plan_type": "premium",
        "subscription_status": "active",
        "subscription_id": payment_id,
    })
    supabase_request("POST", "/payments", {
        "user_id": user_id,
        "paymongo_payment_id": payment_id,
        "paymongo_link_id": link_id,
        "amount": SUPPORTER_AMOUNT,
        "currency": "PHP",
        "status": "paid",
        "paid_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
    })
    logger.info(f"User {user_id} upgraded to Supporter")


def check_payment_link(link_id: str) -> dict:
    req = urllib.request.Request(
        f"https://api.paymongo.com/v1/links/{link_id}",
        headers={"Authorization": paymongo_auth_header()},
        method="GET"
    )
    with urllib.request.urlopen(req, timeout=10) as resp:
        data = json.loads(resp.read().decode())
        attrs = data["data"]["attributes"]
        payments = attrs.get("payments", [])
        return {
            "status": attrs.get("status", "unpaid"),
            "paid": attrs.get("status") == "paid",
            "amount": attrs.get("amount", 0),
            "payment_id": payments[0]["id"] if payments else None,
        }


# ── Vercel WSGI handler ───────────────────────────────────────────────────────

class handler(BaseHTTPRequestHandler):

    def do_OPTIONS(self):
        self.send_response(200)
        self._send_cors_headers()
        self.end_headers()

    def do_GET(self):
        path = urlparse(self.path).path
        query = parse_qs(urlparse(self.path).query)

        if path in ("/api/health", "/api/quiz/health"):
            available = [p["name"] for p in PROVIDERS if os.getenv(p["api_key_env"])]
            self._json_response(200, {"status": "healthy", "providers_configured": available, "version": "1.0.0"})

        elif path == "/api/paymongo/verify":
            # Poll payment status — called from success page
            link_id = query.get("link_id", [None])[0]
            user_id = query.get("user_id", [None])[0]
            if not link_id or not user_id:
                self._json_response(400, {"detail": "link_id and user_id required"})
                return
            try:
                result = check_payment_link(link_id)
                if result["paid"] and result["payment_id"]:
                    mark_user_premium(user_id, result["payment_id"], link_id)
                self._json_response(200, result)
            except Exception as e:
                logger.error(f"Verify failed: {e}")
                self._json_response(500, {"detail": str(e)})

        else:
            self._json_response(404, {"detail": "Not found"})

    def do_POST(self):
        path = urlparse(self.path).path
        if path == "/api/quiz/generate":
            self._handle_generate()
        elif path == "/api/quiz/test-gemini":
            self._handle_test()
        elif path == "/api/paymongo/create-link":
            self._handle_create_link()
        elif path == "/api/paymongo/webhook":
            self._handle_webhook()
        else:
            self._json_response(404, {"detail": "Not found"})

    def _handle_create_link(self):
        try:
            content_length = int(self.headers.get("Content-Length", 0))
            body = json.loads(self.rfile.read(content_length).decode())
            user_id = body.get("userId")
            user_email = body.get("userEmail")

            if not user_id or not user_email:
                self._json_response(400, {"detail": "userId and userEmail required"})
                return
            if not PAYMONGO_SECRET_KEY:
                self._json_response(500, {"detail": "PayMongo not configured on server"})
                return

            result = create_paymongo_link(user_id, user_email)
            self._json_response(200, result)
        except Exception as e:
            logger.error(f"Create link failed: {e}")
            self._json_response(500, {"detail": str(e)})

    def _handle_webhook(self):
        try:
            content_length = int(self.headers.get("Content-Length", 0))
            body = self.rfile.read(content_length)
            sig = self.headers.get("Paymongo-Signature", "")

            if sig and not verify_webhook_signature(body, sig):
                self._json_response(401, {"detail": "Invalid signature"})
                return

            event = json.loads(body.decode())
            event_type = event.get("data", {}).get("attributes", {}).get("type", "")
            logger.info(f"PayMongo webhook: {event_type}")

            if event_type in ("link.payment.paid", "payment.paid"):
                attrs = event["data"]["attributes"]
                payment_data = attrs.get("data", {})
                metadata = (
                    payment_data.get("attributes", {}).get("metadata", {})
                    or attrs.get("metadata", {})
                )
                user_id = metadata.get("user_id")

                if not user_id:
                    remarks = payment_data.get("attributes", {}).get("remarks", "")
                    if "user_id:" in remarks:
                        user_id = remarks.split("user_id:")[1].strip()

                if user_id:
                    payment_id = payment_data.get("id", event["data"]["id"])
                    link_id = attrs.get("id", "")
                    mark_user_premium(user_id, payment_id, link_id)
                    self._json_response(200, {"received": True, "upgraded": True})
                else:
                    logger.warning("Webhook: no user_id found in metadata")
                    self._json_response(200, {"received": True, "upgraded": False})
            else:
                self._json_response(200, {"received": True})

        except Exception as e:
            logger.error(f"Webhook failed: {e}")
            self._json_response(500, {"detail": str(e)})

    def _handle_generate(self):
        try:
            content_type = self.headers.get("Content-Type", "")
            content_length = int(self.headers.get("Content-Length", 0))
            body = self.rfile.read(content_length)

            if "multipart/form-data" not in content_type:
                self._json_response(400, {"detail": "Expected multipart/form-data"})
                return

            boundary = None
            for part in content_type.split(";"):
                part = part.strip()
                if part.startswith("boundary="):
                    boundary = part[9:].strip()
                    break

            if not boundary:
                self._json_response(400, {"detail": "No boundary in multipart data"})
                return

            fields, files = self._parse_multipart(body, boundary.encode())
            if "file" not in files:
                self._json_response(400, {"detail": "No file provided"})
                return

            result = generate_quiz_logic(
                files["file"]["content"],
                files["file"].get("filename", "upload.txt"),
                int(fields.get("question_count", "10")),
                fields.get("difficulty", "medium")
            )
            self._json_response(200, result)

        except ValueError as e:
            self._json_response(400, {"detail": str(e)})
        except Exception as e:
            logger.error(f"Generate failed: {e}")
            self._json_response(500, {"detail": str(e)})

    def _handle_test(self):
        try:
            text, provider = generate_with_fallback('Respond with this exact JSON: [{"status": "ok"}]')
            self._json_response(200, {"success": True, "provider": provider, "response": text})
        except Exception as e:
            self._json_response(500, {"detail": str(e)})

    def _parse_multipart(self, body: bytes, boundary: bytes):
        fields, files = {}, {}
        delimiter = b"--" + boundary
        parts = body.split(delimiter)
        for part in parts[1:]:
            if part in (b"--\r\n", b"--", b"\r\n") or part.startswith(b"--"):
                continue
            if b"\r\n\r\n" not in part:
                continue
            headers_raw, content = part.split(b"\r\n\r\n", 1)
            if content.endswith(b"\r\n"):
                content = content[:-2]
            disposition = {}
            for line in headers_raw.decode("utf-8", errors="ignore").split("\r\n"):
                if "Content-Disposition" in line:
                    for item in line.split(";"):
                        item = item.strip()
                        if "=" in item:
                            k, v = item.split("=", 1)
                            disposition[k.strip()] = v.strip().strip('"')
            name = disposition.get("name", "")
            filename = disposition.get("filename", "")
            if filename:
                files[name] = {"content": content, "filename": filename}
            else:
                fields[name] = content.decode("utf-8", errors="ignore")
        return fields, files

    def _json_response(self, status: int, data: dict):
        body = json.dumps(data).encode()
        self.send_response(status)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(body)))
        self._send_cors_headers()
        self.end_headers()
        self.wfile.write(body)

    def _send_cors_headers(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "GET, POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type, Paymongo-Signature")

    def log_message(self, fmt, *args):
        logger.info(f"{self.address_string()} - {fmt % args}")