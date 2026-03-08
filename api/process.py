import os
import json
import tempfile
import shutil
from typing import Dict, List, Optional, Any
from pathlib import Path
from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
import uvicorn
import pdfplumber
import docx
import io
import re
import random

# Additional imports for enhanced file support
try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

app = FastAPI()

# Configure CORS for Vercel frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["https://examzz.vercel.app", "http://localhost:5173", "http://localhost:4173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Simple data models (without Pydantic)
class ExtractedContent:
    def __init__(self, text: str, headings: List[str] = None, keywords: List[str] = None, sentences: List[str] = None):
        self.text = text
        self.headings = headings or []
        self.keywords = keywords or []
        self.sentences = sentences or []

class Question:
    def __init__(self, question: str, options: List[str], correctAnswer: int, questionType: str, explanation: Optional[str] = None):
        self.question = question
        self.options = options
        self.correctAnswer = correctAnswer
        self.questionType = questionType
        self.explanation = explanation

class QuizRequest:
    def __init__(self, content: Dict[str, Any], questionCount: int = 10, quizType: str = "quiz"):
        self.content = ExtractedContent(**content)
        self.questionCount = questionCount
        self.quizType = quizType

def clean_extracted_text(text: str) -> str:
    """Clean extracted text while preserving structure"""
    text = re.sub(r'[^\x20-\x7E\n\r\t]', ' ', text)
    text = re.sub(r'[ \t]+', ' ', text)  # Only collapse spaces/tabs
    text = re.sub(r'\n{3,}', '\n\n', text)  # Max 2 consecutive newlines
    return text.strip()

def extract_sentences(text: str) -> List[str]:
    """Extract sentences from text, handling bullet points and numbered lists"""
    sentences = []
    
    # Split on newlines first
    lines = text.split('\n')
    
    for line in lines:
        line = line.strip()
        if not line or len(line) < 20:
            continue
            
        # If line is short enough, use directly
        if len(line) <= 500:
            sentences.append(line)
            continue
            
        # If line is too long, split on punctuation
        # Split on sentence-ending punctuation
        parts = re.split(r'([.!?]+\s+)', line)
        current = ''
        
        for part in parts:
            current += part
            if re.search(r'[.!?]\s*$', current) and len(current.strip()) > 20:
                if len(current.strip()) <= 500:
                    sentences.append(current.strip())
                current = ''
        
        # Push any remainder
        if current.strip() and len(current.strip()) > 20:
            sentences.append(current.strip())
    
    # Remove duplicates while preserving order
    seen = set()
    unique_sentences = []
    for sentence in sentences:
        if sentence not in seen:
            seen.add(sentence)
            unique_sentences.append(sentence)
    
    return unique_sentences

def extract_headings(text: str) -> List[str]:
    """Extract headings from text"""
    lines = text.split('\n')
    headings = []
    
    for line in lines:
        line = line.strip()
        if len(line) > 5 and (
            re.match(r'^#+\s+', line) or  # Markdown headers
            re.match(r'^[A-Z][^.:]*:$', line) or  # Title case
            re.match(r'^\d+\.\s+', line)  # Numbered lists
        ):
            headings.append(line)
    
    return headings[:20]  # Limit to 20 headings

def extract_keywords(text: str, headings: List[str]) -> List[str]:
    """Extract keywords from text"""
    # Common words to exclude
    common_words = {
        'this', 'that', 'these', 'those', 'there', 'their', 'the', 'a', 'an',
        'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with',
        'by', 'from', 'is', 'are', 'was', 'were', 'be', 'been', 'have', 'has',
        'had', 'do', 'does', 'did', 'will', 'would', 'could', 'should', 'may',
        'might', 'can', 'shall', 'must', 'what', 'which', 'who', 'when',
        'where', 'why', 'how', 'all', 'any', 'both', 'each', 'few', 'more',
        'most', 'other', 'such', 'only', 'own', 'same', 'so', 'than', 'up', 'out'
    }
    
    # Extract words (3+ characters)
    words = re.findall(r'\b[a-zA-Z]{3,}\b', text.lower())
    
    # Count occurrences
    word_count = {}
    for word in words:
        if word not in common_words:
            word_count[word] = word_count.get(word, 0) + 1
    
    # Sort by frequency and take top 30
    keywords = sorted(word_count.items(), key=lambda x: x[1], reverse=True)[:30]
    return [word for word, count in keywords]

def process_document_content(text: str) -> ExtractedContent:
    """Process document text and extract content"""
    cleaned_text = clean_extracted_text(text)
    sentences = extract_sentences(cleaned_text)
    headings = extract_headings(cleaned_text)
    keywords = extract_keywords(cleaned_text, headings)
    
    return ExtractedContent(
        text=cleaned_text,
        headings=headings,
        keywords=keywords,
        sentences=sentences
    )

def extract_pdf_text(file_content: bytes) -> str:
    """Extract text from PDF using pdfplumber"""
    try:
        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            text = ""
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return ""

def extract_docx_text(file_content: bytes) -> str:
    """Extract text from DOCX using python-docx"""
    try:
        doc = python_docx.Document(io.BytesIO(file_content))
        return "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
    except Exception as e:
        print(f"DOCX extraction error: {e}")
        return ""

def extract_pptx_text(file_content: bytes) -> str:
    """Extract text from PPTX using python-pptx"""
    if not HAS_PPTX:
        print("python-pptx not installed")
        return ""
    try:
        prs = Presentation(io.BytesIO(file_content))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines)
    except Exception as e:
        print(f"PPTX extraction error: {e}")
        return ""

def extract_xlsx_text(file_content: bytes) -> str:
    """Extract text from XLSX using openpyxl"""
    if not HAS_OPENPYXL:
        print("openpyxl not installed")
        return ""
    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_content), read_only=True)
        lines = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = " ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    lines.append(row_text.strip())
        return "\n".join(lines)
    except Exception as e:
        print(f"XLSX extraction error: {e}")
        return ""

def extract_text_from_file(file: UploadFile) -> str:
    """Extract text from uploaded file"""
    content = file.file.read()
    filename = file.filename.lower() if file.filename else ""
    
    if filename.endswith('.pdf'):
        return extract_pdf_text(content)
    elif filename.endswith(('.docx', '.doc')):
        return extract_docx_text(content)
    elif filename.endswith(('.pptx', '.ppt')):
        return extract_pptx_text(content)
    elif filename.endswith(('.xlsx', '.xls')):
        return extract_xlsx_text(content)
    elif filename.endswith('.txt'):
        return content.decode('utf-8', errors='ignore')
    else:
        raise ValueError(f"Unsupported file type: {filename}")

def generate_definition_questions(content: ExtractedContent) -> List[Question]:
    """Generate definition questions with better distractors"""
    questions = []
    definition_patterns = [
        r'(\w+)\s+is\s+(?:defined\s+as\s+)?(?:the\s+)?process\s+of\s+([^,.]+)',
        r'(\w+)\s+is\s+(?:defined\s+as\s+)?(?:a\s+)?([^,.]+)',
        r'(\w+)\s+refers\s+to\s+([^,.]+)',
        r'(\w+)\s+means\s+([^,.]+)',
        r'The\s+(\w+)\s+is\s+(?:the\s+)?([^,.]+)',
    ]
    
    for sentence in content.sentences:
        for pattern in definition_patterns:
            match = re.search(pattern, sentence, re.IGNORECASE)
            if match:
                term = match.group(1).strip()
                definition = match.group(2).strip() if len(match.groups()) > 1 else ""
                
                if term and definition and len(term) > 2:
                    # Generate better distractors from keywords
                    distractors = [k for k in content.keywords 
                                 if k.lower() != term.lower() and len(k) > 3][:3]
                    
                    if len(distractors) >= 2:
                        options = [definition] + distractors[:2]
                        random.shuffle(options)
                        questions.append(Question(
                            question=f"What is {term}?",
                            options=options,
                            correctAnswer=options.index(definition),
                            questionType="definition",
                            explanation=f"According to the text: '{sentence}'"
                        ))
                break  # Only one question per sentence
    
    return questions

def generate_fill_blank_questions(content: ExtractedContent) -> List[Question]:
    """Generate fill-in-the-blank questions with better distractors"""
    questions = []
    
    for sentence in content.sentences:
        # Find important words (keywords or capitalized words)
        important_words = [
            w for w in re.findall(r'\b[a-zA-Z]{4,}\b', sentence)
            if (w.lower() in content.keywords or 
                (w[0].isupper() and len(w) > 4))
        ]
        
        if not important_words:
            continue
            
        word = random.choice(important_words)
        blanked_sentence = re.sub(rf'\b{re.escape(word)}\b', '______', sentence, count=1, flags=re.IGNORECASE)
        
        if blanked_sentence != sentence:
            # Generate distractors from keywords
            distractors = [k for k in content.keywords 
                         if k.lower() != word.lower()][:3]
            
            if len(distractors) >= 2:
                options = [word] + distractors[:2]
                random.shuffle(options)
                questions.append(Question(
                    question=f"Fill in the blank: {blanked_sentence}",
                    options=options,
                    correctAnswer=options.index(word),
                    questionType="fill-blank",
                    explanation=f"The correct answer is '{word}'. Full sentence: '{sentence}'"
                ))
    
    return questions

def generate_keyword_questions(content: ExtractedContent) -> List[Question]:
    """Generate keyword-based questions"""
    questions = []
    
    if len(content.keywords) >= 4:
        # Select a keyword as the "correct" answer
        correct_keyword = content.keywords[0]
        
        # Get 3 other keywords as distractors
        other_keywords = [k for k in content.keywords[1:4] if k.lower() != correct_keyword.lower()]
        
        if len(other_keywords) >= 3:
            options = [correct_keyword] + other_keywords
            questions.append(Question(
                question=f"Which of the following is most important in the context?",
                options=options,
                correctAnswer=0,
                questionType="keyword",
                explanation=f"Based on the text discussion"
            ))
    
    return questions

def generate_multiple_choice_questions(content: ExtractedContent) -> List[Question]:
    """Generate multiple choice questions"""
    questions = []
    
    for sentence in content.sentences[:10]:  # Limit to first 10 sentences
        words = re.findall(r'\b[a-zA-Z]{4,}\b', sentence)
        
        if len(words) >= 5:
            question_word = words[0]
            
            # Get other words as options
            other_words = [w for w in words[1:4] if w.lower() != question_word.lower()]
            
            if len(other_words) >= 3:
                options = [question_word] + other_words
                questions.append(Question(
                    question=f"According to the text, which word is most relevant?",
                    options=options,
                    correctAnswer=0,
                    questionType="multiple-choice",
                    explanation=f"Context from: '{sentence}'"
                ))
    
    return questions

def generate_quiz_from_content(request: QuizRequest) -> List[Question]:
    """Generate quiz questions from content"""
    all_questions = []
    
    # Generate different question types
    definition_questions = generate_definition_questions(request.content)
    fill_blank_questions = generate_fill_blank_questions(request.content)
    keyword_questions = generate_keyword_questions(request.content)
    multiple_choice_questions = generate_multiple_choice_questions(request.content)
    
    all_questions.extend(definition_questions)
    all_questions.extend(fill_blank_questions)
    all_questions.extend(keyword_questions)
    all_questions.extend(multiple_choice_questions)
    
    # Shuffle and limit to requested count
    import random
    random.shuffle(all_questions)
    
    return all_questions[:request.questionCount]

@app.post("/process-document")
async def process_document(file: UploadFile = File(...)):
    """Process uploaded document and extract content"""
    try:
        text = extract_text_from_file(file)
        
        if not text or len(text.strip()) < 50:
            raise HTTPException(status_code=400, detail="Not enough content to process")
        
        content = process_document_content(text)
        
        return {
            "text": content.text,
            "headings": content.headings,
            "keywords": content.keywords,
            "sentences": content.sentences
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Document processing failed: {str(e)}")

@app.post("/generate-quiz")
async def generate_quiz(request: dict):
    """Generate quiz questions from content"""
    try:
        # Manually create QuizRequest from dict
        quiz_request = QuizRequest(
            content=request.get("content", {}),
            questionCount=request.get("questionCount", 10),
            quizType=request.get("quizType", "quiz")
        )
        
        questions = generate_quiz_from_content(quiz_request)
        
        if not questions:
            raise HTTPException(status_code=400, detail="Could not generate questions from content")
        
        # Convert to dict for JSON response
        question_dicts = []
        for q in questions:
            question_dicts.append({
                "question": q.question,
                "options": q.options,
                "correctAnswer": q.correctAnswer,
                "questionType": q.questionType,
                "explanation": q.explanation
            })
        
        return {
            "questions": question_dicts,
            "count": len(question_dicts)
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Quiz generation failed: {str(e)}")

@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {"status": "healthy", "service": "document-processor"}

# For Vercel serverless deployment
handler = app

# For local development
if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
