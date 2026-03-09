import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import openpyxl
import os
import logging
import traceback

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def extract_text_from_file(file_path: str) -> str:
    """
    Extracts clean text from various file formats.
    :param file_path: Path to the uploaded file.
    :return: Extracted text as a string.
    :raises ValueError: If file type is unsupported or extraction fails.
    """
    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    try:
        logger.info(f"Extracting text from file: {file_path} (type: {ext})")
        
        if ext == '.pdf':
            with fitz.open(file_path) as doc:
                for page_num, page in enumerate(doc):
                    page_text = page.get_text()
                    if page_text.strip():
                        text += page_text + "\n"
                    logger.info(f"Processed PDF page {page_num + 1}/{len(doc)}")
                        
        elif ext in ['.docx', '.doc']:
            doc = Document(file_path)
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())
            text = "\n".join(paragraphs)
            logger.info(f"Processed Word document with {len(paragraphs)} paragraphs")
            
        elif ext in ['.pptx', '.ppt']:
            prs = Presentation(file_path)
            slide_count = 0
            for slide_num, slide in enumerate(prs.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                if slide_text.strip():
                    text += slide_text + "\n"
                    slide_count += 1
                logger.info(f"Processed slide {slide_num + 1}/{len(prs.slides)}")
            logger.info(f"Processed presentation with {slide_count} text-containing slides")
            
        elif ext in ['.xlsx', '.xls']:
            wb = openpyxl.load_workbook(file_path, read_only=True)
            sheet_count = 0
            for sheet_name, sheet in wb.items():
                sheet_text = ""
                for row_num, row in enumerate(sheet.iter_rows()):
                    row_values = [str(cell.value) for cell in row if cell.value and str(cell.value).strip()]
                    if row_values:
                        sheet_text += " ".join(row_values) + "\n"
                if sheet_text.strip():
                    text += f"Sheet: {sheet_name}\n{sheet_text}\n"
                    sheet_count += 1
                logger.info(f"Processed sheet '{sheet_name}' with {row_num + 1} rows")
            logger.info(f"Processed Excel workbook with {sheet_count} data-containing sheets")
            
        else:
            raise ValueError(f"Unsupported file type: {ext}")
        
        # Clean up: Remove extra newlines and whitespace
        text = "\n".join([line.strip() for line in text.splitlines() if line.strip()])
        
        if not text:
            raise ValueError(f"No text extracted from file: {file_path}")
            
        logger.info(f"Successfully extracted {len(text)} characters from {file_path}")
        return text
        
    except Exception as e:
        logger.error(f"Extraction failed for {file_path}: {str(e)}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        raise ValueError(f"Extraction failed: {str(e)}")

def validate_file_before_extraction(file_path: str, max_size_mb: int = 10) -> bool:
    """
    Validates file before extraction.
    :param file_path: Path to the file
    :param max_size_mb: Maximum file size in MB
    :return: True if valid, raises ValueError if invalid
    """
    if not os.path.exists(file_path):
        raise ValueError(f"File does not exist: {file_path}")
    
    file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
    if file_size_mb > max_size_mb:
        raise ValueError(f"File size {file_size_mb:.2f}MB exceeds maximum {max_size_mb}MB")
    
    ext = os.path.splitext(file_path)[1].lower()
    supported_extensions = ['.pdf', '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls']
    
    if ext not in supported_extensions:
        raise ValueError(f"Unsupported file type: {ext}. Supported types: {supported_extensions}")
    
    return True
